import re
from pathlib import Path
from .base import BuddyPlugin


class PPTXPlugin(BuddyPlugin):
    name = "pptx"
    supported_extensions = ["pptx", "ppt"]
    description = "Read and edit PowerPoint presentations"

    def execute(self, file_path: str, task: str) -> str:
        path = self.validate_file(file_path)
        try:
            from pptx import Presentation
        except ImportError:
            return "❌ python-pptx not installed. Run: pip install python-pptx"

        lo = task.lower()
        if any(k in lo for k in ["read", "text", "content", "slides", "show"]):
            return self._read(path, Presentation)
        if any(k in lo for k in ["count", "how many"]):
            return self._count(path, Presentation)
        if "add slide" in lo:
            return self._add_slide(path, task, Presentation)
        if any(k in lo for k in ["replace", "find"]):
            return self._replace(path, task, Presentation)
        return self._read(path, Presentation)

    def _read(self, path: Path, Presentation) -> str:
        prs = Presentation(str(path))
        parts = []
        for i, slide in enumerate(prs.slides, 1):
            texts = [s.text.strip() for s in slide.shapes if hasattr(s, "text") and s.text.strip()]
            if texts:
                parts.append(f"[Slide {i}]\n" + "\n".join(texts))
        if not parts:
            return "❌ No text found in presentation."
        full = "\n\n".join(parts)
        if len(full) > 3000:
            return f"📊 Slides (first 3000 chars):\n\n{full[:3000]}…"
        return f"📊 Slides:\n\n{full}"

    def _count(self, path: Path, Presentation) -> str:
        prs = Presentation(str(path))
        n = len(prs.slides)
        return f"📊 {path.name} → {n} slide{'s' if n != 1 else ''}."

    def _add_slide(self, path: Path, task: str, Presentation) -> str:
        m = re.search(r'add slide[:\s]+(.+)', task, re.I)
        if not m:
            return "❌ Format: add slide: Slide Title"
        title = m.group(1).strip()
        prs = Presentation(str(path))
        layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(layout)
        slide.shapes.title.text = title
        prs.save(str(path))
        return f"✅ Slide '{title}' added to {path.name}"

    def _replace(self, path: Path, task: str, Presentation) -> str:
        m = re.search(r'replace\s+"(.+?)"\s+with\s+"(.+?)"', task, re.I)
        if not m:
            m = re.search(r'replace\s+(.+?)\s+with\s+(.+)', task, re.I)
        if not m:
            return '❌ Format: replace "old" with "new"'
        old, new = m.group(1).strip(), m.group(2).strip()
        prs = Presentation(str(path))
        count = 0
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text_frame"):
                    for para in shape.text_frame.paragraphs:
                        for run in para.runs:
                            if old in run.text:
                                run.text = run.text.replace(old, new)
                                count += 1
        prs.save(str(path))
        return f"✅ Replaced {count} instance(s) in {path.name}"
